import streamlit as st
import os
import pandas as pd
from pptx import Presentation
from dotenv import load_dotenv

# --- THE NUCLEAR SSL FIX ---
import ssl
import urllib3

# Disable SSL warnings so they don't clutter your terminal
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# Force Python to ignore certificate validation errors globally
try:
    _create_unverified_https_context = ssl._create_unverified_context
except AttributeError:
    pass
else:
    ssl._create_default_https_context = _create_unverified_https_context

# Force requests/httpx to ignore SSL (Common for VPNs/Zscaler)
os.environ["CURL_CA_BUNDLE"] = ""
os.environ["REQUESTS_CA_BUNDLE"] = ""
# ---------------------------

# The Modern, LCEL-based LangChain imports
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.documents import Document
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser

# 1. Load your API key from the .env file
load_dotenv()

# 2. Setup Page configuration
st.set_page_config(page_title="Data Chat Agent", layout="wide")

# Inject Neo-Brutalist CSS for a modern, minimal UI
st.markdown("""
    <style>
    .stApp {
        background-color: #f4f4f0;
        color: #000000;
        font-family: 'Courier New', Courier, monospace;
    }
    .stButton>button {
        border: 3px solid #000000;
        border-radius: 0px;
        background-color: #ffde59;
        color: #000000;
        font-weight: bold;
        box-shadow: 4px 4px 0px #000000;
        transition: all 0.1s ease;
    }
    .stButton>button:active {
        box-shadow: 0px 0px 0px #000000;
        transform: translate(4px, 4px);
    }
    .stTextInput>div>div>input {
        border: 3px solid #000000;
        border-radius: 0px;
    }
    .stFileUploader>div>div>button {
        border: 2px solid #000000;
        border-radius: 0px;
    }
    </style>
    """, unsafe_allow_html=True)

st.title("Eizou")

# 3. Helper Functions to Read Files
def read_excel(file):
    """Reads Excel and converts rows to a readable markdown table."""
    df = pd.read_excel(file)
    return df.to_markdown()

def read_ppt(file):
    """Extracts text from all shapes in a PowerPoint slide."""
    prs = Presentation(file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

# 4. Process and Vectorize Documents
def process_documents(uploaded_files):
    documents = []
    
    for file in uploaded_files:
        file_extension = file.name.split('.')[-1].lower()
        raw_text = ""
        
        if file_extension in ['xlsx', 'xls']:
            raw_text = read_excel(file)
        elif file_extension in ['pptx', 'ppt']:
            raw_text = read_ppt(file)
            
        if raw_text:
            documents.append(Document(page_content=raw_text, metadata={"source": file.name}))
            
    if not documents:
        return None

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=2000) # Determines chunk size and overlap between chunks
    chunks = text_splitter.split_documents(documents)
    
    embeddings = GoogleGenerativeAIEmbeddings(model="gemini-embedding-001", transport="rest")
    vector_store = FAISS.from_documents(chunks, embeddings)
    return vector_store

# 5. Sidebar Interface for Uploads
with st.sidebar:
    st.header("📂 Data Sources")
    uploaded_files = st.file_uploader(
        "Upload Excel or PPT files", 
        type=["xlsx", "xls", "pptx", "ppt"], 
        accept_multiple_files=True
    )
    
    if st.button("Process Data"):
        if uploaded_files:
            with st.spinner("Crunching data..."):
                vector_store = process_documents(uploaded_files)
                if vector_store:
                    st.session_state.vector_store = vector_store
                    st.success("Data processed and saved to memory!")
                else:
                    st.error("Could not extract text from the uploaded files.")
        else:
            st.warning("Please upload files first.")

# 6. Chat Interface
if "messages" not in st.session_state:
    st.session_state.messages = []

for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

# Helper function for the new LCEL pipeline
def format_docs(docs):
    return "\n\n".join(doc.page_content for doc in docs)

if prompt := st.chat_input("Ask a question about your uploaded data..."):
    st.chat_message("user").markdown(prompt)
    st.session_state.messages.append({"role": "user", "content": prompt})

    if "vector_store" in st.session_state:
        llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0, transport="rest")
        
        system_prompt = (
            "You are a meticulous enterprise data assistant. Use the provided context to answer the user's question. "
            "CRITICAL INSTRUCTION: When the user asks you to 'list', 'group', or 'show all' items, you MUST output the COMPLETE and exhaustive list from the context. "
            "Do NOT summarize, do NOT truncate, and do NOT say 'here are a few examples'. Extract and provide every single matching item. "
            "If you don't know the answer based on the context, say you don't know.\n\n"
            "Context: {context}"
        )
        prompt_template = ChatPromptTemplate.from_messages([
            ("system", system_prompt),
            ("human", "{input}"),
        ])
        
        retriever = st.session_state.vector_store.as_retriever(search_kwargs={"k": 50}) # Determines how many chunks are retrieved at once
        
        # The Modern LCEL Pipeline (Replaces the broken .chains module)
        rag_chain = (
            {"context": retriever | format_docs, "input": RunnablePassthrough()}
            | prompt_template
            | llm
            | StrOutputParser()
        )
        
        with st.spinner("Thinking..."):
            # Invoke the LCEL chain directly
            answer = rag_chain.invoke(prompt)
            
        with st.chat_message("assistant"):
            st.markdown(answer)
        st.session_state.messages.append({"role": "assistant", "content": answer})
    else:
        st.error("Please upload and process some files in the sidebar first!")