from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Depends, Header
from fastapi.middleware.cors import CORSMiddleware
import os
import pandas as pd
from pptx import Presentation
from dotenv import load_dotenv, find_dotenv
import io
from supabase import create_client, Client
import PyPDF2
import boto3
import re
import json
import random

# --- LANGCHAIN IMPORTS ---
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_aws import ChatBedrockConverse, BedrockEmbeddings
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.documents import Document
from langchain_core.output_parsers import StrOutputParser

load_dotenv(find_dotenv())

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- SUPABASE INIT ---
url: str = os.getenv("SUPABASE_URL")
key: str = os.getenv("SUPABASE_KEY")
supabase: Client = create_client(url, key)

# --- AUTH & RBAC DEPENDENCIES ---
def get_current_user(authorization: str = Header(None)):
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing Token")
    
    token = authorization.split(" ")[1]
    
    try:
        temp_client = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_KEY"))
        user_response = temp_client.auth.get_user(token)
        user_id = user_response.user.id
        email = user_response.user.email
        
        profile = supabase.table("profiles").select("role").eq("id", user_id).execute()
        role = profile.data[0]["role"] if profile.data else "user"
        
        return {"id": user_id, "email": email, "role": role}
    except Exception as e:
        print(f"Auth verification error: {str(e)}")
        raise HTTPException(status_code=401, detail="Invalid or Expired Token")

def require_admin(user: dict = Depends(get_current_user)):
    if user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")
    return user

# --- AUTH ENDPOINTS ---
@app.post("/login")
def login(email: str = Form(...), password: str = Form(...)):
    try:
        temp_client = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_KEY"))
        res = temp_client.auth.sign_in_with_password({"email": email, "password": password})
        token = res.session.access_token
        
        profile = supabase.table("profiles").select("role").eq("id", res.user.id).execute()
        role = profile.data[0]["role"] if profile.data else "user"
        
        return {"token": token, "role": role, "email": res.user.email}
    except Exception:
        raise HTTPException(status_code=401, detail="Invalid credentials")

# --- ADMIN ENDPOINTS ---
@app.post("/admin/users")
def create_user(email: str = Form(...), password: str = Form(...), role: str = Form(...), admin: dict = Depends(require_admin)):
    try:
        admin_client = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_KEY"))
        
        res = admin_client.auth.admin.create_user({
            "email": email,
            "password": password,
            "email_confirm": True
        })
        new_user_id = res.user.id
        
        admin_client.table("profiles").update({"role": role}).eq("id", new_user_id).execute()
        return {"status": "User provisioned successfully"}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/admin/projects")
def create_project(name: str = Form(...), admin: dict = Depends(require_admin)):
    res = supabase.table("projects").insert({"name": name}).execute()
    return {"status": "Project created", "data": res.data}

@app.post("/admin/assign")
def assign_user(user_email: str = Form(...), project_id: int = Form(...), admin: dict = Depends(require_admin)):
    profile = supabase.table("profiles").select("id").eq("email", user_email).execute()
    if not profile.data:
        raise HTTPException(status_code=404, detail="User not found")
    
    user_id = profile.data[0]["id"]
    try:
        supabase.table("project_users").insert({"project_id": project_id, "user_id": user_id}).execute()
        return {"status": "User assigned"}
    except Exception:
        return {"status": "User already assigned"}

@app.get("/admin/list_users")
def list_users(admin: dict = Depends(require_admin)):
    profiles = supabase.table("profiles").select("*").execute()
    return {"users": profiles.data}

@app.get("/admin/files")
def list_files(admin: dict = Depends(require_admin)):
    try:
        files = supabase.table("project_files").select("*").execute()
        return {"files": files.data}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/admin/projects/{project_id}")
def delete_project(project_id: int, admin: dict = Depends(require_admin)):
    try:
        supabase.table("projects").delete().eq("id", project_id).execute()
        return {"status": "Project deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/admin/projects/{project_id}")
def edit_project(project_id: int, name: str = Form(...), admin: dict = Depends(require_admin)):
    try:
        supabase.table("projects").update({"name": name}).eq("id", project_id).execute()
        return {"status": "Project updated"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/admin/users/{user_id}")
def delete_user(user_id: str, admin: dict = Depends(require_admin)):
    try:
        supabase.auth.admin.delete_user(user_id)
        supabase.table("profiles").delete().eq("id", user_id).execute()
        return {"status": "User deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/admin/users/{user_id}/role")
def update_user_role(user_id: str, role: str = Form(...), admin: dict = Depends(require_admin)):
    try:
        supabase.table("profiles").update({"role": role}).eq("id", user_id).execute()
        return {"status": "Role updated"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/admin/files/{file_id}")
def delete_file(file_id: int, admin: dict = Depends(require_admin)):
    try:
        file_record = supabase.table("project_files").select("*").eq("id", file_id).execute()
        if file_record.data:
            path = f"project_{file_record.data[0]['project_id']}/{file_record.data[0]['file_name']}"
            supabase.storage.from_("project_files").remove([path])
        
        supabase.table("project_files").delete().eq("id", file_id).execute()
        return {"status": "File deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# --- USER ENDPOINTS ---
@app.get("/projects")
def get_user_projects(user: dict = Depends(get_current_user)):
    try:
        if user["role"] == "admin":
            projects = supabase.table("projects").select("*").execute()
            return {"projects": projects.data}
        else:
            assignments = supabase.table("project_users").select("project_id").eq("user_id", user["id"]).execute()
            
            if not assignments.data:
                return {"projects": []}
                
            project_ids = [str(a["project_id"]) for a in assignments.data]
            projects = supabase.table("projects").select("*").in_("id", project_ids).execute()
            
            return {"projects": projects.data}
            
    except Exception as e:
        print(f"❌ Error fetching projects: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Database error: {str(e)}")

# --- AI DATA ENDPOINTS ---
@app.post("/upload")
async def upload_files(files: list[UploadFile] = File(...), project_id: int = Form(...), model: str = Form(...), admin: dict = Depends(require_admin)):
    documents = []
    
    for file in files:
        contents = await file.read()
        ext = file.filename.split('.')[-1].lower()
        file_path = f"project_{project_id}/{file.filename}"
        
        file_url = "#" 
        
        try:
            supabase.storage.from_("project_files").upload(file_path, contents, file_options={"upsert": "true"})
            file_url = supabase.storage.from_("project_files").get_public_url(file_path)
            
            supabase.table("project_files").insert({
                "project_id": project_id,
                "file_name": file.filename,
                "file_url": file_url
            }).execute()
        except Exception as e:
            print(f"⚠️ Storage Error for {file.filename}: {str(e)}")

        raw_text = ""
        try:
            if ext in ['xlsx', 'xls']:
                raw_text = pd.read_excel(io.BytesIO(contents)).to_markdown()
            elif ext == 'csv':
                raw_text = pd.read_csv(io.BytesIO(contents)).to_markdown()
            elif ext in ['pptx', 'ppt']:
                prs = Presentation(io.BytesIO(contents))
                raw_text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            elif ext == 'pdf':
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(contents))
                raw_text = "\n".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])
            else:
                raw_text = contents.decode('utf-8', errors='ignore')
        except Exception as e:
            print(f"⚠️ Parsing Error for {file.filename}: {str(e)}")

        if raw_text.strip():
            documents.append(Document(
                page_content=raw_text, 
                metadata={
                    "source": file.filename, 
                    "project_id": project_id, 
                    "file_url": file_url
                }
            ))
            
    if not documents:
        raise HTTPException(status_code=400, detail="No readable text found in the uploaded files.")

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=2000, chunk_overlap=200)
    chunks = text_splitter.split_documents(documents)
    
    if "gemini" in model.lower():
        embeddings = GoogleGenerativeAIEmbeddings(
            model="gemini-embedding-001", 
            google_api_key=os.getenv("GOOGLE_API_KEY"), 
            transport="rest"
        )
    else:
        bedrock_client = boto3.client(
            service_name="bedrock-runtime",
            region_name=os.getenv("AWS_DEFAULT_REGION", "us-east-1")
        )
        embeddings = BedrockEmbeddings(client=bedrock_client, model_id="amazon.titan-embed-text-v2:0")
        
    try:
        texts = [chunk.page_content for chunk in chunks]
        metadatas = [chunk.metadata for chunk in chunks]
        
        raw_vectors = embeddings.embed_documents(texts)
        truncated_vectors = [vec[:1024] for vec in raw_vectors]
        
        records = []
        for text, meta, vec in zip(texts, metadatas, truncated_vectors):
            records.append({
                "content": text,
                "metadata": meta,
                "embedding": vec
            })
            
        supabase.table("project_documents").insert(records).execute()
        
    except Exception as e:
        print(f"❌ Vector DB Error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to sync with AI memory: {str(e)}")
    
    return {"status": "success", "message": f"Knowledge base synced. {len(files)} files processed."}

@app.post("/chat")
async def chat(message: str = Form(...), project_id: int = Form(...), model: str = Form(...), user: dict = Depends(get_current_user)):
    
    if "gemini" in model.lower():
        embeddings = GoogleGenerativeAIEmbeddings(model="gemini-embedding-001", google_api_key=os.getenv("GOOGLE_API_KEY"), transport="rest")
        llm = ChatGoogleGenerativeAI(model=model, google_api_key=os.getenv("GOOGLE_API_KEY"), temperature=0, transport="rest")
    else:
        bedrock_client = boto3.client(
            service_name="bedrock-runtime",
            region_name=os.getenv("AWS_DEFAULT_REGION", "ap-south-1"),
            aws_access_key_id=os.getenv("AWS_ACCESS_KEY_ID"),
            aws_secret_access_key=os.getenv("AWS_SECRET_ACCESS_KEY")
        )
        embeddings = BedrockEmbeddings(client=bedrock_client, model_id="amazon.titan-embed-text-v2:0")
        llm = ChatBedrockConverse(client=bedrock_client, model_id="openai.gpt-oss-20b-1:0", temperature=0)
    
    try:
        # Embed the raw message directly
        query_embedding = embeddings.embed_query(message)[:1024]
        
        # Only filter by project_id, let Titan embeddings do the semantic matching
        rpc_response = supabase.rpc("match_project_documents", {
            "query_embedding": query_embedding,
            "match_count": 20,
            "filter": {"project_id": int(project_id)} 
        }).execute()
        
        chunks = rpc_response.data or []
        context_text = ""
        potential_filenames = set()

        for row in chunks:
            meta = row.get("metadata", {})
            name = meta.get("source", "Unknown")
            content = row.get("content", "")
            context_text += f"\n---\nDOCUMENT: {name}\nCONTENT: {content}\n"
            potential_filenames.add(name)

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Search error: {str(e)}")
    
    system_prompt = (
        "You are an enterprise data assistant. Answer based strictly on the context provided.\n\n"
        "### OUTPUT RULES:\n"
        "1. VISUALS: If the user asks for a graph, trend, or comparison, you MUST generate a valid JSON object "
        "wrapped in a ```chart ... ``` block. Use 'bar', 'line', or 'pie'.\n"
        "Structure: {{ \"type\": \"bar\", \"data\": {{ \"labels\": [\"Jan\", \"Feb\"], \"datasets\": [{{ \"label\": \"Revenue\", \"data\": [100, 200], \"backgroundColor\": \"#0A56D0\" }}] }}, \"options\": {{ \"responsive\": true }} }}\n\n"
        "2. CITATIONS: At the very end of your response, after any text or charts, you MUST list the document names used.\n"
        "Format: SOURCES: [File1.pdf, File2.xlsx]\n\n"
        "Context: {context}"
    )

    prompt_template = ChatPromptTemplate.from_messages([
        ("system", system_prompt), ("human", "{input}"),
    ])
    
    rag_chain = prompt_template | llm | StrOutputParser()
    
    try:
        full_response = rag_chain.invoke({"context": context_text, "input": message})
        
        cited_sources = []
        clean_answer = full_response
        
        if "SOURCES:" in full_response:
            parts = full_response.rsplit("SOURCES:", 1)
            clean_answer = parts[0].strip()
            
            raw_names = re.findall(r"\[(.*?)\]", parts[1])
            
            if raw_names:
                names_list = [n.strip() for n in raw_names[0].split(",")]
                
                file_data = supabase.table("project_files") \
                    .select("file_name, file_url") \
                    .in_("file_name", names_list) \
                    .eq("project_id", project_id) \
                    .execute()
                
                url_map = {f["file_name"]: f["file_url"] for f in file_data.data}
                
                for name in names_list:
                    cited_sources.append({
                        "name": name, 
                        "url": url_map.get(name, "#")
                    })

        return {"answer": clean_answer, "sources": cited_sources}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/projects/{project_id}/recommendations")
async def get_recommendations(project_id: int, user: dict = Depends(get_current_user)):
    try:
        # Fetch chunks to build context
        docs = supabase.table("project_documents").select("content, metadata").contains("metadata", {"project_id": project_id}).limit(40).execute()
        
        if not docs.data:
            return {"questions": []}
            
        sample_size = min(5, len(docs.data))
        random_chunks = random.sample(docs.data, sample_size)
        
        context = ""
        for d in random_chunks:
            src = d.get("metadata", {}).get("source", "Unknown")
            context += f"--- Source Document: {src} ---\n{d['content']}\n\n"
        
        # Request a simple array of strings back
        prompt = (
            "You are an analytical AI. Read the text chunks below.\n"
            "Extract 3 to 5 highly specific facts.\n"
            "Turn each fact into a specific question that can be answered purely by the text.\n"
            "Return ONLY a valid JSON array of strings. Do not use markdown tags.\n"
            "Example: [\"What is the Q3 revenue?\", \"How many Anvils and Rams are present in the Clutch Press shop?\"]\n\n"
            f"Context:\n{context[:4000]}"
        )
        
        bedrock_client = boto3.client(
            service_name="bedrock-runtime",
            region_name=os.getenv("AWS_DEFAULT_REGION", "ap-south-1"),
            aws_access_key_id=os.getenv("AWS_ACCESS_KEY_ID"),
            aws_secret_access_key=os.getenv("AWS_SECRET_ACCESS_KEY")
        )
        llm = ChatBedrockConverse(client=bedrock_client, model_id="openai.gpt-oss-20b-1:0", temperature=0.7)
        
        response = llm.invoke(prompt)
        
        raw_content = response.content
        text_content = "".join(b.get("text", "") for b in raw_content if isinstance(b, dict)) if isinstance(raw_content, list) else str(raw_content)
        
        clean_text = text_content.replace("```json", "").replace("```", "").strip()
        questions_array = json.loads(clean_text)
        
        # Ensure we only send back strings
        valid_questions = [str(q) for q in questions_array if isinstance(q, str)]
        
        return {"questions": valid_questions}
        
    except Exception as e:
        print(f"⚠️ Recommendation generation failed: {str(e)}")
        return {"questions": []}